import io
import json
import logging
import os.path
from datetime import datetime
from typing import List
from urllib.parse import urljoin

import dash
import requests
from argon2 import PasswordHasher
from bs4 import BeautifulSoup
from dash import dcc
from pptx import Presentation

title = "Voreinteilungen 👀"

with open("config.json", "r") as f:
    config = json.load(f)

hasher = PasswordHasher()
logging.basicConfig(level=logging.ERROR, format="%(levelname)s: %(message)s")

def validate_template_structure(config):
    logging.info("Starting template structure validation.")

    # Check if "template" exists
    if "template" not in config:
        logging.error("'template' key is missing in config.")
        return False
    logging.info("'template' key found.")

    # Check for required keys in 'template'
    required_keys = ["path", "id_template_ref-team", "id_template_ref-single",
                     "template_ref-team_mapping", "template_ref-single_mapping"]

    for key in required_keys:
        if key not in config["template"]:
            logging.error(f"Required key '{key}' is missing in 'template' config.")
            return False
        logging.info(f"'{key}' key found in 'template' config.")

    # Check if 'league_mapping' exists (optional)
    if "league_mapping" in config["template"]:
        logging.info("'league_mapping' key found (optional).")
    else:
        logging.info("'league_mapping' key is missing but optional.")

    # Validate the 'path' contains a valid .pptx file
    path = config["template"]["path"]
    if not os.path.exists(path):
        logging.error(f"Path '{path}' does not exist.")
        return False
    if not os.path.isfile(path):
        logging.error(f"Path '{path}' is not a file.")
        return False
    if not path.lower().endswith(".pptx"):
        logging.error(f"Path '{path}' is not a PowerPoint (.pptx) file.")
        return False

    logging.info(f"PowerPoint (.pptx) file found at '{path}'.")

    logging.info("Template structure validation successful. Download is enabled.")
    return True
template = validate_template_structure(config)


def update_config():
    with open("config.json", "w") as f:
        json.dump(config, f)


def get_password_hash_for_user(username: str) -> str:
    if username not in config["auth"]:
        return ""
    return config["auth"][username]["password"]


def set_password_hash_for_user(user: str, pw_hash: str) -> None:
    config["auth"][user]["password"] = pw_hash
    update_config()


def url_builder(users: List[List[str]], prefix="refs") -> str:
    args = f"?{prefix}=" + "_".join(users[0])
    for user in users[1:]:
        args += f"&{prefix}=" + "_".join(user)
    return args


base_url = "https://www.dfbnet.org"
dfbnet_landing = "https://www.dfbnet.org/spielplus/login.do"
dfbnet_login = "https://www.dfbnet.org/spielplus/oauth/login"
search = "https://www.dfbnet.org/sria/mod_sria/offenespielelist.do?reqCode=view"


def get_grouped_users(user_groups):
    group_links = {}  # name: list_of_users
    for key in config["grouped_users"]:
        if config["grouped_users"][key]["group"] not in user_groups:
            continue
        group_links[key] = config["grouped_users"][key]["users"]
    return group_links


def get_single_users(user_groups):
    single_user_links = []  # list_of_users
    for group in user_groups:
        if not isinstance(group, list) and len(group) != 2:
            continue
        single_user_links += [group]
    return single_user_links


def get_ref_req(vorname, nachname, datedelta):
    return {"staffel": "", "msa_id": "0", "status": "4", "date": datetime.today().strftime("%d.%m.%Y"),
            "datedelta": str(datedelta), "srvorname": vorname, "srnachname": nachname, "spieltag": ""}


def search_link(web_content, keyword):
    for el in web_content.find_all("a"):
        if el.text == keyword:
            return el["href"]
    return None


class Ref:
    def __init__(self, ref_args: List[str]):
        self.role = ref_args[0][0]
        self.name = ref_args[0][1]
        self.state = ref_args[1]
        self.atspl = ref_args[0][2]

    def __repr__(self):
        return f"{self.role}: {self.name} ({self.state})"

    def __eq__(self, other):
        return self.role == other.role and self.name == other.name and self.state == other.state

    def __hash__(self):
        return hash((self.role, self.name, self.state))


class Match:
    def __init__(self, match_args: List[str], ref_state):
        split_args = match_args[1].split("\n")
        if len(split_args) == 3:
            self.date = " ".join(split_args[1:])
        elif len(split_args) == 2:
            self.date = " ".join(split_args)
        else:
            self.date = " ".join(split_args)
        self.date = datetime.strptime(self.date, "%d.%m.%Y %H:%M")
        split_args = match_args[2].split("\n")
        if len(split_args) == 2:
            self.staffel, self.match_id = split_args
        else:
            self.staffel, self.match_id = " ".join(split_args), ""
        split_args = match_args[4].split("\n")
        if len(split_args) == 2:
            self.home, self.location = split_args
        else:
            self.home, self.location = split_args[0], ""
        self.guest = match_args[5]

        valid_roles = ["SR", "SRA1", "SRA2", "BEO", "PA", "4OF"]
        search_for_name = False
        team_args = []
        current_ref = []
        current_ATSPL = ""
        for a in match_args[7].split("\n"):
            if "ATS" in a or a == "":
                continue
            elif "-->" in a:
                current_ATSPL = a.replace("-->", "").strip()
                continue
            elif a in valid_roles:
                if not search_for_name:
                    search_for_name = True
                else:
                    current_ref.append("")
                    current_ref.append(current_ATSPL)
                    team_args.append(current_ref)
                current_ref = [a]
                current_ATSPL = ""
                continue
            else:
                current_ref.append(a)
                current_ref.append(current_ATSPL)
                team_args.append(current_ref)
            current_ref = []
            search_for_name = False
        if search_for_name:
            current_ref.append("")
            current_ref.append(current_ATSPL)
            team_args.append(current_ref)

        self.team = [Ref(args) for args in zip(team_args, ref_state)]

    def __repr__(self):
        return f"{self.date}\n{self.staffel}\n{self.home} v. {self.guest}\n{self.location}\n{self.team}"

    def __eq__(self, other):
        return (self.date == other.date and
                self.staffel == other.staffel and
                self.match_id == other.match_id and
                self.home == other.home and
                self.location == other.location and
                self.guest == other.guest and
                self.team == other.team)

    def __hash__(self):
        return hash((self.date, self.staffel, self.home, self.location, self.guest, self.team))

    def create_powerpoint_output(self):
        ref = None
        sra1 = None
        sra2 = None

        for r in self.team:
            if not r.name:
                continue
            name_split = r.name.split()
            name = " ".join(name_split[:-1]) + "\n" + name_split[-1]
            if r.role == "SR":
                ref = name
            elif r.role == "SRA1":
                sra1 = name
            elif r.role == "SRA2":
                sra2 = name
            else:
                pass

        if "league_mapping" in config["template"] and self.staffel in config["template"]["league_mapping"]:
            staffel_name = config["template"]["league_mapping"][self.staffel].upper()
        else:
            staffel_name = self.staffel

        return_values = [
            self.date.strftime("%d.%m.%Y"),
            self.date.strftime("%H:%M"),
            self.location,
            self.home.upper(),
            self.guest.upper(),
            staffel_name,
            ref,
        ]
        if sra1 or sra2:
            return_values += [sra1, sra2]

        return return_values


def parse_icons(contents):
    rows = contents.find_all("tr")
    return_list = []

    for row in rows:
        icons = row.find_all("img")
        if len(icons) == 0:
            if "ATS" not in row.text:
                return_list.append("")
        else:
            for icon in icons:
                match icon["alt"]:
                    case "Ansetzung bestätigt.":
                        return_list += ["✅"]
                    case "Ansetzung nicht bestätigt.":
                        return_list += ["❓"]
                    case "Vorläufige Einteilung":
                        return_list += ["✘"]
    return return_list


def parse_matches(web_page):
    match = None
    try:
        matches = []

        table = web_page.find("table", attrs={"class": "sportView"})
        rows = table.find_all("tr", recursive=False)[1:]
        for row in rows:
            elements = []
            match = row.find_all("td", recursive=False)
            for el in match:
                if el.get_text() == "Keine Einträge gefunden!":
                    return []
                elements += [el.get_text("\n").strip().replace("\xa0", "")]
            matches += [Match(elements, parse_icons(match[-2]))]
    except Exception as e:
        print(e)
        print(match)
    return matches


def prepare_search_session(username, password):
    s = requests.Session()
    s.get(dfbnet_landing)
    resp = s.get(dfbnet_login)
    x = BeautifulSoup(resp.text, "html.parser")
    auth_webpage = x.find(id="kc-form-login")["action"]
    resp = s.post(auth_webpage, data={
        "username": username,
        "password": password,
        "credentialId": "",
    })
    x = BeautifulSoup(resp.text, "html.parser")
    new_link = search_link(x, "Schiriansetzung")
    resp = s.get(urljoin(base_url, new_link))
    x = BeautifulSoup(resp.text, "html.parser")
    new_link = search_link(x, "Ansetzung")
    s.get(urljoin(base_url, new_link))
    return s


def search_ref(session, nachname, vorname):
    resp = session.post(search, data=get_ref_req(nachname=nachname, vorname=vorname, datedelta=999))
    x = BeautifulSoup(resp.text, "html.parser")
    return parse_matches(x)


def create_instagram_template(data):
    if not template:
        return dash.no_update
    prs = Presentation(config["template"]["path"])
    layout_3_refs = config["template"]["id_template_ref-team"]
    layout_1_refs = config["template"]["id_template_ref-single"]

    lookup_table_3 = {x: i for i, x in enumerate(config["template"]["template_ref-team_mapping"].values())}
    lookup_table_1 = {x: i for i, x in enumerate(config["template"]["template_ref-single_mapping"].values())}

    for match in data:
        if len(match) == 9:
            # 3 refs
            lookup_table = lookup_table_3
            slide_layout = prs.slide_layouts[layout_3_refs]
        elif len(match) == 7:
            # 1 refs
            lookup_table = lookup_table_1
            slide_layout = prs.slide_layouts[layout_1_refs]
        else:
            logging.error("Invalid input!")
            return dash.no_update
        slide = prs.slides.add_slide(slide_layout)
        for i, shape in enumerate(slide.placeholders):
            if i not in lookup_table:
                continue
            shape.text = match[lookup_table[i]]
    io_buffer = io.BytesIO()
    prs.save(io_buffer)
    return dcc.send_bytes(io_buffer.getvalue(), "matchday.pptx")